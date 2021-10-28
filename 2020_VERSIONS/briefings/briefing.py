##################################################
####### IMPACTS FORECAST BRIEFING GENERATOR ######
##################################################
#
### THIS PROGRAM GENERATES A POWERPOINT CONTAINING
### OBSERVATIONAL AND MODEL DATA FOR THE PURPOSES
### OF ORGANIZING A FORECAST BRIEFING. THIS
### PROGRAM CALLS A 'NAMELIST' FILE CONTAINING
### VALUES TO KEY VARIABLES THAT THE PROGRAM USES.
#
### FOR ISSUES OR ADDITIONAL HELP, CONTACT:
###   Joseph Finlon
###   jfinlon@uw.edu
#
# REQUIRED FILES: 'namelist.txt'
# REQUIRED PYTHON PACKAGES: pptx, re, urllib, numpy, pandas
#
##################################################
### IMPORT PYTHON PACKAGES
from datetime import time as dtime
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
#from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from PIL import Image
import os, sys, re, glob
import urllib
import numpy as np
import pandas as pd
##################################################
### AUXILLARY FUNCTIONS
# Get timestamp of most recent regional radar image
def get_latest_radar(radarID):
    url = 'http://weather.rap.ucar.edu/radar/nids/images/K' + radarID + '/'

    with urllib.request.urlopen(url) as response:
        encoding = response.info().get_param('charset', 'utf8')
        html = response.read().decode(encoding)

    links = re.findall(r'href=[\'"]?([^\'" >]+)', html) # Find all url instances on the parent (radar site) webpage
    rad_datestr = links[-1][:15] # The last link on the webpage is the latest image...parse the timestamp

    return (url, rad_datestr)

def dt_nearestTimes(yr, mn, dy, hr):
    # Gets the datetime corresponding to the nearest 1, 6, and 12 hour time for model initialization, etc.
    if hr<5:
        past12hr = datetime(yr, mn, dy-1, 12, 0)
        nearest12hr = datetime(yr, mn, dy-1, 12, 0)
        nearest6hr = datetime(yr, mn, dy-1, 18, 0)
        if hr>=1:
            current12hr = datetime(yr, mn, dy, 0, 0)
        else:
            current12hr = datetime(yr, mn, dy-1, 12, 0)
    elif hr>=5 and hr<10:
        past12hr = datetime(yr, mn, dy, 0, 0)
        current12hr = past12hr
        nearest12hr = datetime(yr, mn, dy, 0, 0)
        nearest6hr = datetime(yr, mn, dy, 0, 0)
    elif hr>=10 and hr<16:
        past12hr = datetime(yr, mn, dy, 0, 0)
        nearest12hr = datetime(yr, mn, dy, 0, 0)
        nearest6hr = datetime(yr, mn, dy, 6, 0)
        if hr>=13:
            current12hr = datetime(yr, mn, dy, 12, 0)
        else:
            current12hr = nearest12hr
    elif hr>=16 and hr<20:
        past12hr = datetime(yr, mn, dy, 12, 0)
        current12hr = past12hr
        nearest12hr = datetime(yr, mn, dy, 12, 0)
        nearest6hr = datetime(yr, mn, dy, 12, 0)
    else:
        past12hr = datetime(yr, mn, dy, 12, 0)
        current12hr = past12hr
        nearest12hr = datetime(yr, mn, dy, 12, 0)
        nearest6hr = datetime(yr, mn, dy, 18, 0)
        
    if hr<2:
        nearest2hr = datetime(yr, mn, dy-1, hr-2, 0)
    else:
        nearest2hr = datetime(yr, mn, dy, hr-2, 0)
    
    nearest1hr = datetime(yr, mn, dy, hr, 0)
    
    return (past12hr, current12hr, nearest1hr, nearest2hr, nearest6hr, nearest12hr)

def dt2string(dt_object, init_time):
    dt_utc_full = datetime.strftime(dt_object, '%Y%m%d%H%M')
    dt_utc_short = datetime.strftime(dt_object, '%H UTC %m/%d')
    dt_local = datetime.strftime(dt_object - timedelta(hours=5), '%a %-I%p %m/%d')
    elapsed_time = int((dt_object - init_time).total_seconds() / 3600)
    
    return (dt_utc_full, dt_utc_short, dt_local, elapsed_time)

def get_latest_image(product):
    remotefile, localfile, prod_str, utc_str, local_str, elapsed = img_paths[product]
    print(remotefile)
    not_found = False
    if remotefile==localfile: # verify that a non-URL image is on the local server
        if os.path.isfile(remotefile) is False:
            not_found = True
    elif remotefile is None:
        not_found = True
    else:
        if os.path.exists(localfile):
            os.system('rm -f {:s}'.format(localfile))
            
        try:
            urllib.request.urlretrieve(remotefile, localfile)
        except AttributeError:
            urllib.urlretrieve(remotefile, localfile)
        except:
            not_found = True
            
    if not_found:
        print('  WARNING: Did not find {:s} image'.format(product))
        return (None, prod_str, utc_str, local_str, elapsed)
    else: # Return the path and other info
        return (localfile, prod_str, utc_str, local_str, elapsed)
    
def get_gpm_overpasses(nearest1hr, region='east'):
    '''
    Gets the first six GPM overpasses from the briefing time in the IMPACTS domain (east or west)
    region: 'east' or 'west'
    '''
    gpm_dir = '/home/disk/funnel/impacts/archive/ops/gpm'
    if os.path.isdir(gpm_dir): # only run if this path can be found
        currDate_string = datetime.strftime(nearest1hr, '%Y%m%d')
        currDate_int = int(currDate_string) # will use this to incriment dates by 1
        dates = [str(currDate_int), str(currDate_int+1), str(currDate_int+2)]
        files = []
        for date in dates:
            files.append(glob.glob(gpm_dir + '/' + date + '/ops.gps.' + date + '*.overpass_' + region + '.png'))
        files_flattened = [val for sublist in files for val in sublist]

        for filenum, file in enumerate(files_flattened[:6]): # loop through first 6 files to construct overpass dict
            product_name = 'overpass' + str(filenum + 1)
            remote_file = file
            local_file = remote_file
            timestamp = file.split('.')[2] # get timestamp from filename
            gpm_dt = datetime.strptime(timestamp[:10], '%Y%m%d%H')
            gpm_fullStr, gpm_shortStr, gpm_locStr, gpm_delta = dt2string(gpm_dt, present_time)
            product_string = '{} Overpass'.format(gpm_shortStr)
            img_paths[product_name] = (remote_file, local_file, product_string,
                                       fcst_shortStr, fcst_locStr, fcst_delta)

# def get_latest_image(product, useTime=None, resize=False):
# 	"""
# 	Function to get the image from a given product and catalog the time it's valid
# 	useTime: None or datetime that the image is valid
# 	resize: True will cut the image width and height in half to reduce file size

# 	Returns a tuple with:

# 	(string that is the full address of the image)

# 	The product name MUST HAVE an entry in the img_paths dictionary,
# 	otherwise None is returned
# 	"""

# 	if product not in img_paths.keys():
# 		print('  Unable to find path to product {}'.format(product))
# 		return None

# 	if useTime is not None:
# 		validTime = useTime
# 	else: # Assign current time for the product
# 		validTime = datetime.utcnow()

# 	# Parse out the info (path is the remote URL, ext is what the saved image will be named, product string)
# 	path, ext, prod_str = img_paths[product]
# 	origpath = path
# 	not_found = False

# 	if ext is None: # No filename given to save the image
# 		recent_file = path.split('/')[-1] # Just use the filename on the website
		
# 		if os.path.exists(recent_file):
# 			os.system('rm -f {:s}'.format(recent_file))

# 		try:
# 			urllib.request.urlretrieve(path, recent_file)
# 		except AttributeError:
# 			urllib.urlretrieve(path, recent_file)
# 		except:
# 			not_found = True
# 	else:
# 		if os.path.exists(ext):
# 			os.system('rm -f {:s}'.format(ext))
# 		try:
# 			urllib.request.urlretrieve(path, ext)
# 		except AttributeError:
# 			urllib.urlretrieve(path, ext)
# 		except:
# 			not_found = True
# 		recent_file = ext
# 	print(path, recent_file)

# 	if not_found:
# 		print('  WARNING: Did not find {:s} image'.format(product))
# 		return (None, None)
# 	else: # Return the path
# 		if resize is True:
# 			imObj = Image.open(recent_file)
# 			imWidth, imHeight = (imObj.width // 2, imObj.height // 2)
# 			imObj_resized = imObj.resize((imWidth, imHeight))
# 			imObj_resized.save(recent_file)
# 		if path == '':
# 			return (recent_file, validTime)
# 		else:
# 			# return ('/'.join((path,recent_file)), validTime)
# 			return (recent_file, validTime)

def add_timeline(slide, curday):
    # Add the timeline to the top of the slide
    shapes = slide.shapes
    # Set colors to use
    daycolor = {-1:RGBColor(182, 66, 206),
                0: RGBColor(178, 34, 34),
                1: RGBColor(218, 165, 32),
                2: RGBColor(46, 139, 87),
                3: RGBColor(65, 105, 225)}

    # First, past weather (Day -1)
    left = Inches(0)
    top= Inches(0)
    height = Inches(0.4)
    width = Inches(2)
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON,\
        left, top, width, height)
    if (type(curday) is int) and (curday not in [0,1,2,3,4,5]):
        curday = -1
    elif (type(curday) is int) and (curday >= 3):
        curday = 3

    shape.line.color.rgb = daycolor[-1]
    
    # Now figure out the shading
    shape.fill.solid()
    if (type(curday) is list) and (curday[0]==-1):
        shape.fill.fore_color.rgb = daycolor[-1]
    if (type(curday) is int) and (curday == -1):
        shape.fill.fore_color.rgb = daycolor[-1]
    else:
        shape.fill.fore_color.rgb = RGBColor(200,200,200)

    r=shape.text_frame.paragraphs[0].add_run()
    r.text = 'Day -1'
    r.font.size = Pt(14)

    left += width - Inches(0.05)
    width = Inches(2)
    
    # Now for rest of the days
    for day in range(0,4):

        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.line.color.rgb = daycolor[day]
        r = shape.text_frame.paragraphs[0].add_run()
        if day == 3:
            r.text = 'Day {:d}+ ({:%d %b}-)'.format(day, present_time + timedelta(days=day))
        else:
            r.text = 'Day {:d} ({:%d %b})'.format(day, present_time + timedelta(days=day))
        r.font.size = Pt(14)
        left += width - Inches(0.05)
        shape.fill.solid()
        if type(curday) is list:
            filled = False
            for val in curday:
                if val == day:
                    shape.fill.fore_color.rgb = daycolor[val]
                    filled = True
                if filled is False:
                    shape.fill.fore_color.rgb = RGBColor(200,200,200)
        elif curday == day:
            shape.fill.fore_color.rgb = daycolor[curday]
        else:
            shape.fill.fore_color.rgb = RGBColor(200,200,200)

    return slide

def timeline_slide(prs, curday, presentationType, products=None):
    daycolor = {-1:RGBColor(182, 66, 206), 0: RGBColor(178, 34, 34), 1: RGBColor(218, 165, 32),
                2: RGBColor(46, 139, 87), 3: RGBColor(65, 105, 225)}
    slide_layout = prs.slide_layouts[layout['Title Alone']]
    slide = prs.slides.add_slide(slide_layout)
    chevron_left = Inches(0.25); chevron_top= Inches(5.5); chevron_height = Inches(0.4);
    
    if curday==[0, 1, 2]:
        slide.shapes.title.text = 'Day 0-2 Forecast Overview'
        
        # Add graphics
        images = []
        for panelnum in range(2):
            results = get_latest_image(products[panelnum])
            if results[0] is None:
                images.append([])
            else:
                images.append(results)

        lefts = [0.5, 5.5]; tops = [1.25, 1.25];
        widthval = 4.5 ; heightval = 3.7;

        # Now the panels
        for panelnum, panel in enumerate(images):
            if panel == []: # No image found for current product, model, time, and scope configuration
                continue
            try:
                pic = slide.shapes.add_picture(panel[0], left=Inches(lefts[panelnum]), top=Inches(tops[panelnum]),
                                               width=Inches(widthval), height=Inches(heightval))
            except:
                pass
            
        # Add the zoomed-in timeline
        if presentationType=='morning':
            chevron_width = [1.9, 3.8, 3.8]
        else: # evening briefing
            chevron_width = [1.0, 4.25, 4.25]
        
        shapes = slide.shapes
        for day in range(0, 3):
            shape = shapes.add_shape(MSO_SHAPE.CHEVRON, chevron_left, chevron_top, Inches(chevron_width[day]),
                                     chevron_height)
            shape.line.color.rgb = daycolor[day]
            r = shape.text_frame.paragraphs[0].add_run()
            r.text = 'Day {:d} ({:%d %b})'.format(day, present_time + timedelta(days=day)); r.font.size = Pt(14);
            chevron_left += Inches(chevron_width[day]) - Inches(0.05)
            shape.fill.solid(); shape.fill.fore_color.rgb = daycolor[day];
            
        # Add the timeline
        slide = add_timeline(slide, [0, 1, 2])
        return prs
    else: # Day 3+
        slide.shapes.title.text = 'Day 3-5 Forecast Overview'
        
        # Add the zoomed-in timeline
        chevron_width = Inches(9.5/3.)
        shapes = slide.shapes
        for day in range(3, 6):
            shape = shapes.add_shape(MSO_SHAPE.CHEVRON, chevron_left, chevron_top, chevron_width, chevron_height)
            shape.line.color.rgb = daycolor[3]
            r = shape.text_frame.paragraphs[0].add_run()
            r.text = 'Day {:d} ({:%d %b})'.format(day, present_time + timedelta(days=day)); r.font.size = Pt(14);
            chevron_left += chevron_width - Inches(0.05)
            shape.fill.solid(); shape.fill.fore_color.rgb = daycolor[3];
            
        # Add the timeline
        slide = add_timeline(slide, 3)
        return prs
    
def probability_table(prs, titletxt, curday):
    """
    Make an empty slide and insert the probabilities for various events
    """
    slide_layout = prs.slide_layouts[layout['Title Alone']]
    
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titletxt

    rows = 6; cols = 5;
    width = Inches(9); height = Inches(4);
    top = Inches(2); left = Inches(0.5);
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.first_col = True; table.first_row = True; table.horz_banding = True;
    
    
    ## Write in the table text
    
    # Col titles
    c=table.cell(0,0).text_frame.paragraphs[0].text = 'Objective'
    table.cell(0,1).text_frame.paragraphs[0].text = '12Z - 15Z'
    table.cell(0,2).text_frame.paragraphs[0].text = '15Z - 18Z'
    table.cell(0,3).text_frame.paragraphs[0].text = '18Z - 21Z'
    table.cell(0,4).text_frame.paragraphs[0].text = '21Z - 11Z'
    
    # Row titles
    table.cell(1,0).text_frame.paragraphs[0].text = 'Convection Initiation'
    table.cell(2,0).text_frame.paragraphs[0].text = 'Severe Weather'
    table.cell(3,0).text_frame.paragraphs[0].text = 'Upscale Growth'
    table.cell(4,0).text_frame.paragraphs[0].text = 'Backbuilding'
    table.cell(5,0).text_frame.paragraphs[0].text = 'Hydromet.'

    # Vertically center everything
    for y in range(rows):
        for x in range(cols):
            c=table.cell(y,x)
            c.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            cf=c.text_frame
            cf.paragraphs[0].alignment = PP_ALIGN.CENTER
            if x > 0 and y > 0:
                cf.paragraphs[0].text = '-'

    # Add the timeline
    slide = add_timeline(slide,curday)

    return prs

def timeline_table(prs, titletxt, curday, time_periods):
    """
    Make an empty slide and insert a timeline for the day
    Forecaster will enter notes
    'time_periods': list of 2-char strings for the UTC times,
    beginning and ends included
    """
    
    slide_layout = prs.slide_layouts[layout['Title Alone']]
    
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = titletxt

    # Insert a line object
    shapes = slide.shapes
    width = Inches(9); height = Inches(0.1);
    top = Inches(3); left = Inches(0.5);
    tline = shapes.add_shape(MSO_SHAPE.RECTANGLE,left,top,width,height)
    tline.fill.solid()
    tline.fill.fore_color.rgb = RGBColor(178,34,34)
    
    # Insert bracket objects
    width = Inches(0.5); height = Inches(2.);
    top = Inches(2.5); left = Inches(1.25);
    times = ['{:s}-{:s} UTC'.format(time_periods[x], time_periods[x+1]) for x in range(len(time_periods)-1)]
    
    for b in np.arange(len(times)):
        brace = shapes.add_shape(MSO_SHAPE.LEFT_BRACE,left+Inches(0.25+b*2.2),top,width,height)
        brace.fill.background()
        brace.line.color.rgb = RGBColor(0,0,0) 
        brace.line.width = Pt(6)
        brace.rotation = -90.0
    
        txt         =   slide.shapes.add_textbox(Inches(0.75)+Inches(b*2.2),top+Inches(1),height,width)
        tf          =   txt.text_frame
        para        =   tf.add_paragraph()
        r           =   para.add_run()
        r.text      =   times[b]
        r.font.size = Pt(24)
        r.font.bold = True
        para.alignment = PP_ALIGN.CENTER
    
    # Add the timeline
    slide = add_timeline(slide,curday)
    
    return prs
    
def bumper_slide(prs, title, curday, start_time, end_time):
    # Choose a blank slide
    slide_layout = prs.slide_layouts[layout['Segue']]
    
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    
    # Change the title
    slide.shapes.title.text = title
    
    # List the date and time in the subtitle
    if 'Current Weather' in title:
        slide.placeholders[1].text = "{:%H UTC %d %b %Y}".format(start_time)
    elif 'Past 24' in title:
        slide.placeholders[1].text = "{:%H UTC %d %b %Y} through {:%H UTC %d %b %Y}".format(start_time, end_time)
    elif 'Day 0' not in title:
        slide.placeholders[1].text = "{:%H UTC %d %b %Y} through {:%H UTC %d %b %Y}".format(start_time, end_time)
    else:
        slide.placeholders[1].text = "Now through {:%H UTC %d %b %Y}".format(end_time)

    # Add timeline
    dayposs = re.search('Day (\d)', title)
    add_timeline(slide, curday)

    return prs
    
def objectives_slide(prs, title, curday):
    # Choose a blank slide
    slide_layout = prs.slide_layouts[layout['Bullet Slide']]
    
    # Add the slide to the presentation
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    # Add timeline
    dayposs = re.search('Day (\d)', title)
    add_timeline(slide, curday)

    return prs

def full_summary(prs, title, curday, valid=None):
    slide_layout = prs.slide_layouts[layout['Bullet Slide']]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # Add timeline
    dayposs = re.search('Day (\d)', title)

    if title not in ['Discussion Summary', 'Forecast Timeline', 'Briefing Overview']:
        add_timeline(slide, curday)

    return prs
    
def full_slide_image(prs, product, curday, imgTime, ftime=None, width=None, height=None, link=False, title=None):
    # Take "product" and make a full-slide image with title out of it
    # Grab the latest image
    results = get_latest_image(product)
#     try:
#         results = get_latest_image(product, imgTime)
#     except IOError:
#         results = [None]

    # Get a blank slide layout and add it to the presentation
    if title is not None:
        slide_layout = prs.slide_layouts[layout['Title Alone']]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
    else:
        slide_layout = prs.slide_layouts[layout['Blank Slide']]
        slide = prs.slides.add_slide(slide_layout)
    
#     if (width is not None) or (height is not None):
# #         title.top=Inches(0.5)
#         title.left = Inches(0)
#         title.width=Inches(10)
#     else:
#         title.top=Inches(3)
#         title.left = Inches(7.5)
#         title.width=Inches(2.0)
    
#     p = title.text_frame.paragraphs[0]
#     r = p.add_run()
#     r.text = img_paths[product][2]
#     r.font.size=Pt(32)
    
    if link:
        hlink = r.hyperlink
        hlink.address = link
    
    if ftime is not None:
        d = p.add_run()
        d.text = '\n\n' + ftime.strftime('%d %b %H UTC')
        d.font.size = Pt(28)

    if results[0] is None:
         # Didn't find the image.  Create slide anyway and add the timeline.
        dayposs = re.search('Day (\d)', product)
        add_timeline(slide, curday)
        
        return prs
    
    # Add the image
    if height is not None:
        left_balanced = (10-width)/2.
        try:
            pic = slide.shapes.add_picture(results[0], left=Inches(left_balanced), top=Inches(0.5),
                                           width=Inches(width), height=Inches(height))
        except IOError:
            pic = None
    elif width is not None:
        left_balanced = (10-width)/2.
        try:
            pic = slide.shapes.add_picture(results[0], left=Inches(left_balanced), top=Inches(0.5), width=Inches(width))
        except IOError:
            pic = None
    else:
        try:
            pic = slide.shapes.add_picture(results[0], left=Inches(0), top=Inches(0.5), width=Inches(7))
        except IOError:
            pic = None
    
    # Add timeline
    dayposs = re.search('Day (\d)', product)
    add_timeline(slide, curday)
    
    return prs
    
def two_panel_image(prs, products, curday, imgTimes, link=False, lowertext=False, title=None):
    # Function to plot two images side-by-side
    # Width is set in the function -- not set as an option to the user

    # Get a blank slide layout and add it to the presentation
    slide_layout = prs.slide_layouts[layout['Title Alone']]
    slide = prs.slides.add_slide(slide_layout)

    # Grab the images
    pp = 0

    for product in products:
        widthval = 5
        if pp == 0:
            leftval = 0
        else:
            leftval = widthval

        if title is not None:
            slide.shapes.title.text = title
        else:
            slide.shapes.title.text = ''
#         title = slide.shapes.title; title.top = Inches(0.5); title.left = Inches(leftval); title.width=Inches(widthval);
#         p = title.text_frame.paragraphs[0]
#         r = p.add_run()
#         r.text = ''
#         r.font.size=Pt(40)
        
        txt         =   slide.shapes.add_textbox(Inches(leftval),Inches(1.0),Inches(widthval),Inches(0.5))
        tf          =   txt.text_frame
        para        =   tf.add_paragraph()
        r           =   para.add_run()
        r.text      =   img_paths[product][2]
        r.font.size = Pt(32)
        r.font.bold = True
        para.alignment = PP_ALIGN.CENTER          

        if lowertext != False:
            p = title.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = ''
            r.font.size=Pt(40)
            txt         =   slide.shapes.add_textbox(Inches(leftval),Inches(6.5),Inches(widthval),Inches(0.5))
            tf          =   txt.text_frame
            para        =   tf.add_paragraph()
            r           =   para.add_run()
            r.text      =   lowertext[pp]
            r.font.size = Pt(18)
            r.font.bold = True
            para.alignment = PP_ALIGN.CENTER   

        results = get_latest_image(product)
# 		try:
# 			results = get_latest_image(product)
# 		except IOError:
# 			results = [None]

        if results[0] is None: # Didn't find the image. Create slide anyway and add the timeline.
            pp += 1
            continue
#             dayposs = re.search('Day (\d)', product)
            
#             if dayposs is None:
#                 curday = 0
#             else:
#                 curday = int(dayposs.groups()[0])

#             add_timeline(slide, curday)
            
#             return prs

        # Add the image
        try:
            pic = slide.shapes.add_picture(results[0], left=Inches(leftval), top=Inches(2.0), width=Inches(widthval))
        except IOError:
            pic = None

        pp += 1
    
    # Add timeline
#     dayposs = re.search('Day (\d)', product)
    add_timeline(slide, curday)

    return prs

# def four_panel_image(prs, product_name, curday, day_start, imgTimes, link=None, resize=False):
# 	'''
# 	Download all four images
# 	imgTimes are passed in from the 'times' element in the 'products' structure
# 	'''

# 	images = []
# 	#print(product_name)
# 	pdt = product_name.split('_')
# 	if len(pdt)==5: # For products with an underscore in their name
# 		pdt[0] = pdt[0] + '_' + pdt[1]
# 		pdt.remove(pdt[1])
# 	elif len(pdt)==6: # For products with an underscore in their name
# 		pdt[0] = pdt[0] + '_' + pdt[1] + '_' + pdt[2]
# 		pdt.remove(pdt[1]); pdt.remove(pdt[1]);
# 	elif len(pdt)==7: # For products with an underscore in their name
# 		pdt[0] = pdt[0] + '_' + pdt[1] + '_' + pdt[2] + '_' + pdt[3]
# 		pdt.remove(pdt[1]); pdt.remove(pdt[1]); pdt.remove(pdt[1]);
	
# 	if len(imgTimes)==2:
# 		panelList = ['A','B']
# 		time1 = day_start.replace(hour=0) + timedelta(hours=imgTimes[0])
# 		time2 = day_start.replace(hour=0) + timedelta(hours=imgTimes[1])
# 		times_combined = [time1, time2] # To be used in 'panel' loop below
# 		times = [datetime.strftime(time1, '%H UTC %m/%d'), datetime.strftime(time2, '%H UTC %m/%d')]
# 	elif len(imgTimes)==3:
# 		panelList = ['A','B','C']
# 		time1 = day_start.replace(hour=0) + timedelta(hours=imgTimes[0])
# 		time2 = day_start.replace(hour=0) + timedelta(hours=imgTimes[1])
# 		time3 = day_start.replace(hour=0) + timedelta(hours=imgTimes[2])
# 		times_combined = [time1, time2, time3] # To be used in 'panel' loop below
# 		times = [datetime.strftime(time1, '%H UTC %m/%d'), datetime.strftime(time2, '%H UTC %m/%d'), datetime.strftime(time3, '%H UTC %m/%d')]
# 	else: # 4 panels
# 		panelList = ['A','B','C','D']
# 		time1 = day_start.replace(hour=0) + timedelta(hours=imgTimes[0])
# 		time2 = day_start.replace(hour=0) + timedelta(hours=imgTimes[1])
# 		time3 = day_start.replace(hour=0) + timedelta(hours=imgTimes[2])
# 		time4 = day_start.replace(hour=0) + timedelta(hours=imgTimes[3])
# 		times_combined = [time1, time2, time3, time4] # To be used in 'panel' loop below
# 		times = [datetime.strftime(time1, '%H UTC %m/%d'), datetime.strftime(time2, '%H UTC %m/%d'), datetime.strftime(time3, '%H UTC %m/%d'), datetime.strftime(time4, '%H UTC %m/%d')]

# 	for panelnum, panel in enumerate(panelList):
# 		if resize is True:
# 			results = get_latest_image(pdt[0] + '_' + pdt[1]  + panel + '_' + pdt[2] + '_' + pdt[3], useTime=times_combined[panelnum], resize=True)
# 		else:
# 			results = get_latest_image(pdt[0] + '_' + pdt[1]  + panel + '_' + pdt[2] + '_' + pdt[3], useTime=times_combined[panelnum])
		
# 		if results[0] is None:
# 			images.append([])
# 		else:
# 			images.append(results)

# 	# Get a blank slide
# 	slide_layout = prs.slide_layouts[layout['Title Alone']]
# 	slide = prs.slides.add_slide(slide_layout)
# 	title = slide.shapes.title

# 	lefts = [0, 5, 0, 5]; tops = [0.85, 0.85, 4.2, 4.2];
# 	widthval = 4.5 ; heightval = 3.0;
# 	text_lefts = [2.5,7.5,2.5,7.5]; text_tops = [3.5,3.5,6.8,6.8];
	
# 	# Slide title
# 	txt = slide.shapes.add_textbox(left=Inches(0.0), top=Inches(0.1), width=Inches(10.0), height=Inches(0.5))
# 	tf = txt.text_frame
# 	para = tf.add_paragraph()
# 	r = para.add_run()
# 	r.text = img_paths[pdt[0]+'_'+pdt[1]+'A_'+pdt[2]+'_' +pdt[3]][2] # General product name for the 4-panel plot (grab first time frame for forecast day)
# 	r.font.bold = True
# 	r.font.size = Pt(24)
# 	para.alignment = PP_ALIGN.CENTER
# 	if link is not None:
# 		hlink = r.hyperlink
# 		hlink.address = link
    	
    	
#     # Now the panels
# 	for pnum, p in enumerate(images):
# 		#print(p[0], p[1])
# 		if p == []: # No image found for current product, model, time, and scope configuration
# 			continue
# 		try:
# 			pic = slide.shapes.add_picture(p[0], left=Inches(lefts[pnum]), top=Inches(tops[pnum]), width=Inches(widthval),height=Inches(heightval))
# 		except:
# 			pass

# 		txt = slide.shapes.add_textbox(left=Inches(text_lefts[pnum]), top=Inches(text_tops[pnum]),width=Inches(3.5), height=Inches(0.5))
# 		tf = txt.text_frame
# 		para = tf.add_paragraph()
# 		r = para.add_run()
# 		r.text = times[pnum]
# 		r.font.bold = True
# 		r.font.size = Pt(24)
# 		if link is not None:
# 			hlink = r.hyperlink
# 			hlink.address = link

# 	add_timeline(slide, curday)

# 	return prs

def four_panel_image(prs, products, curday):
    images = []
    for panelnum in range(len(products)):
        results = get_latest_image(products[panelnum])
        images.append(results) # localfile, prod_str, utc_str, local_str, elapsed
#         if results[0] is None:
#             images.append([])
#         else:
#             images.append(results) # localfile, prod_str, utc_str, local_str, elapsed

    # Get a blank slide
    slide_layout = prs.slide_layouts[layout['Blank Slide']]
    slide = prs.slides.add_slide(slide_layout)
#     title = slide.shapes.title

    lefts = [0, 5, 0, 5]; tops = [0.85, 0.85, 4.2, 4.2];
    widthval = 4.5 ; heightval = 3.0;
    text_lefts = [0, 5, 0, 5]; text_tops = [3.5,3.5,6.8,6.8];
#     text_lefts = [2.5,7.5,2.5,7.5]; text_tops = [3.5,3.5,6.8,6.8];
    
    # Slide titles
    txt = slide.shapes.add_textbox(left=Inches(0.0), top=Inches(0.1), width=Inches(3.33), height=Inches(0.5))
    tf = txt.text_frame; para = tf.add_paragraph();
    if images[0][4] < 0:
        r = para.add_run(); r.text = 'Briefing - {} hr'.format(str(-1*images[0][4]));
    else:
        r = para.add_run(); r.text = 'Briefing + {} hr'.format(str(images[0][4]));
    r.font.bold = True; r.font.size = Pt(24);
    para.alignment = PP_ALIGN.LEFT
    
    txt = slide.shapes.add_textbox(left=Inches(3.33), top=Inches(0.1), width=Inches(3.34), height=Inches(0.5))
    tf = txt.text_frame; para = tf.add_paragraph();
    r = para.add_run(); r.text = images[0][2]; r.font.bold = True; r.font.size = Pt(24);
    para.alignment = PP_ALIGN.CENTER
    
    txt = slide.shapes.add_textbox(left=Inches(6.67), top=Inches(0.1), width=Inches(3.33), height=Inches(0.5))
    tf = txt.text_frame; para = tf.add_paragraph();
    r = para.add_run(); r.text = images[0][3]; r.font.bold = True; r.font.size = Pt(24);
    para.alignment = PP_ALIGN.RIGHT
    
    # Now the panels
    for panelnum, panel in enumerate(images):
        txt = slide.shapes.add_textbox(left=Inches(text_lefts[panelnum]), top=Inches(text_tops[panelnum]),
                                       width=Inches(3.5), height=Inches(0.5))
        tf = txt.text_frame
        para = tf.add_paragraph()
        r = para.add_run(); r.text = panel[1]; r.font.bold = True; r.font.size = Pt(18);
        
        if panel[0] is None: # No image found for current product, model, time, and scope configuration
            continue
        try:
            pic = slide.shapes.add_picture(panel[0], left=Inches(lefts[panelnum]), top=Inches(tops[panelnum]),
                                           width=Inches(widthval),height=Inches(heightval))
        except:
            pass

    add_timeline(slide, curday)

    return prs

def six_panel_image(prs, products, curday, plotTimebar=True):
    images = []
    for panelnum in range(len(products)):
        results = get_latest_image(products[panelnum])
        images.append(results)
#         if results[0] is None:
#             images.append([])
#         else:
#             images.append(results)
            
    # Get a blank slide
    slide_layout = prs.slide_layouts[layout['Blank Slide']]
    slide = prs.slides.add_slide(slide_layout)
#     title = slide.shapes.title

    lefts = [0, 3.33, 6.67, 0, 3.33, 6.67]; tops = [0.85, 0.85, 0.85, 4.45, 4.45, 4.45];
    widthval = 3.25 ; heightval = 2.75;
    text_lefts = [0, 3.33, 6.67, 0, 3.33, 6.67]; text_tops = [3.25, 3.25, 3.25, 6.8, 6.8, 6.8];
#     text_lefts = [2.5,7.5,2.5,7.5]; text_tops = [3.5,3.5,6.8,6.8];
    
    # Slide titles
    if plotTimebar is True:
        txt = slide.shapes.add_textbox(left=Inches(0.0), top=Inches(0.1), width=Inches(3.33), height=Inches(0.5))
        tf = txt.text_frame; para = tf.add_paragraph();

        if images[0][4] < 0:
            r = para.add_run(); r.text = 'Briefing - {} hr'.format(str(-1*images[0][4]));
        else:
            r = para.add_run(); r.text = 'Briefing + {} hr'.format(str(images[0][4]));
        r.font.bold = True; r.font.size = Pt(24);
        para.alignment = PP_ALIGN.LEFT

        txt = slide.shapes.add_textbox(left=Inches(3.33), top=Inches(0.1), width=Inches(3.34), height=Inches(0.5))
        tf = txt.text_frame; para = tf.add_paragraph();
        r = para.add_run(); r.text = images[0][2]; r.font.bold = True; r.font.size = Pt(24);
        para.alignment = PP_ALIGN.CENTER

        txt = slide.shapes.add_textbox(left=Inches(6.67), top=Inches(0.1), width=Inches(3.33), height=Inches(0.5))
        tf = txt.text_frame; para = tf.add_paragraph();
        r = para.add_run(); r.text = images[0][3]; r.font.bold = True; r.font.size = Pt(24);
        para.alignment = PP_ALIGN.RIGHT
    
    # Now the panels
    for panelnum, panel in enumerate(images):
        txt = slide.shapes.add_textbox(left=Inches(text_lefts[panelnum]), top=Inches(text_tops[panelnum]),
                                       width=Inches(3.5), height=Inches(0.5))
        tf = txt.text_frame
        para = tf.add_paragraph()
        r = para.add_run(); r.text = panel[1]; r.font.bold = True; r.font.size = Pt(18);
        
        if panel[0] is None: # No image found for current product, model, time, and scope configuration
            continue
        try:
            pic = slide.shapes.add_picture(panel[0], left=Inches(lefts[panelnum]), top=Inches(tops[panelnum]),
                                           width=Inches(widthval),height=Inches(heightval))
        except:
            pass

    add_timeline(slide, curday)

    return prs
##################################################
### PRINT THE PROGRAM PREAMBLE
print('##################################################')
print('#### IMPACTS FORECAST PRESENTATION GENERATOR #####')
print('##################################################')
print('- This code generates a .PPT file of observation')
print('and model fields for use in the IMPACTS briefings.')
print('- Code adapted from Robert Conrick, Univ. Washington')
print('- For questions or assistance, please contact:\n     Joseph Finlon (jfinlon@uw.edu)')
print('##################################################')
##################################################
### READ THE NAMELIST
namelistFile = sys.argv[1] # Gather the namelist file path specified when briefing.py was run
try:
    df = pd.read_csv(namelistFile, sep='=', header=None, names=['Variable','Value'], skipinitialspace=True, comment='#')
    df['Variable'] = df['Variable'].replace(' ','',regex=True) # Remove white space if needed
    df['Value'] = df['Value'].replace('\t','',regex=True) # Remove tabs if needed
except:
    print('Something went wrong parsing the namelist. Make sure it is formatted properly.')

downloadPath = df.loc[0, 'Value'] # The directory where graphics will be downloaded
presentationPath = df.loc[1, 'Value'] # The directory to save the presentation
presentationType = df.loc[2, 'Value'] # The type of briefing to be given (morning, evening)
modelList = df.loc[3, 'Value'].split(',') # Support multiple models for plotting
region = df.loc[5, 'Value']
xSection = df.loc[6, 'Value']
show_weather = df.loc[7, 'Value']
show_shortTerm = df.loc[8, 'Value']
show_detailed = df.loc[9, 'Value']
show_longTerm = df.loc[10, 'Value']
briefingUpdate = df.loc[11, 'Value']
##################################################
### INITIALIZE THE PROGRAM BASED ON THE NAMELIST VALUES
if df.loc[4, 'Value']=='now':
    utcnow = datetime.utcnow()
else:
    utcnow = datetime.strptime(df.loc[4, 'Value'], '%Y%m%d%H')
    print('WARNING: Some model data may have been scrubbed for the selected period...beware.')

# If statements regaring the presentation type and time the program is executed
if presentationType=='morning':
    briefingString = 'morning' # for file naming conventions
    present_time = datetime(utcnow.year, utcnow.month, utcnow.day, 14, 0)
    if utcnow.hour<5 or utcnow.hour>17:
        print('Are you sure you meant to select the morning briefing period? Changing to the evening period now...')
        presentationType = 'evening'

if presentationType=='evening':
    briefingString = 'update' # for file naming conventions
    present_time = datetime(utcnow.year, utcnow.month ,utcnow.day, 23, 0)
    
# Assign model initialization times and sounding, past 1, and 12 hour observation times
past12hr, current12hr, nearest1hr, nearest2hr, nearest6hr, nearest12hr = dt_nearestTimes(
    utcnow.year, utcnow.month, utcnow.day, utcnow.hour)
past12hr_fullStr, past12hr_shortStr, past12hr_locStr, past12hr_delta = dt2string(past12hr, present_time) # Satellite, Analy.
current12hr_fullStr, current12hr_shortStr, current12hr_locStr, current12hr_delta = dt2string(current12hr, present_time)
nearest1hr_fullStr, nearest1hr_shortStr, nearest1hr_locStr, nearest1hr_delta = dt2string(nearest1hr, present_time) # METAR
nearest2hr_fullStr, nearest2hr_shortStr, nearest2hr_locStr, nearest2hr_delta = dt2string(nearest2hr, present_time) # HRRR
nearest6hr_fullStr, nearest6hr_shortStr, nearest6hr_locStr, nearest6hr_delta = dt2string(nearest6hr, present_time) # NAM/GFS
nearest12hr_fullStr, nearest12hr_shortStr, nearest12hr_locStr, nearest12hr_delta = dt2string(nearest12hr, present_time) # WRF
current_fullStr, current_shortStr, current_locStr, current_delta = dt2string(
    datetime(utcnow.year, utcnow.month, utcnow.day, utcnow.hour, utcnow.minute), present_time)
# print(past12hr, current12hr, nearest1hr, nearest2hr, nearest6hr, nearest12hr)
# Assign the cross section orientation to use
xSection_dict = {'n-s':{'value':'C'}, 'nw-se':{'value':'A'}, 'w-e':{'value':'D'}, 'sw-ne':{'value':'C'}}
xSection_value = xSection_dict[xSection]['value']

# Assign the zoomed in region to use
region_dict = {'usne':{'value':'Northeast', 'radString':'northeast', 'metarString':'alb', 'radarLoc':'OKX',
                       'soundLocs':['OKX','PIT','ALB','BUF'], 'plumeLocs':['OKX','CHH','ALB','BUF']},
               'usma':{'value':'Mid-Atlantic', 'radString':'northeast', 'metarString':'bwi', 'radarLoc':'LWX',
                       'soundLocs':['WAL','IAD','RNK','PIT'], 'plumeLocs':['WAL','IAD','RNK','PIT']},
               'usov':{'value':'Ohio Valley', 'radString':'centgrtlakes', 'metarString':'evv', 'radarLoc':'ILN',
                       'soundLocs':['ILN','BNA','ILX','SGF'], 'plumeLocs':['ILN','BNA','ILX','SGF']},
               'usmw':{'value':'Midwest', 'radString':'centgrtlakes', 'metarString':'dtw', 'radarLoc':'LOT',
                       'soundLocs':['APX','DTX','GRB','DVN'], 'plumeLocs':['APX','DTX','GRB','DVN']}
              }
region_value = region_dict[region]['value']
region_radar = region_dict[region]['radString']
region_metar = region_dict[region]['metarString']
# if latestRadarLoc=='default':
#     latestRadarLoc = region_dict[region]['radarLoc']
# if latestSoundLocs[0]=='default':
#     latestSoundLocs = region_dict[region]['soundLocs']
# if fcstSoundLoc=='default':
#     fcstSoundLoc = latestSoundLocs[0]

# Assign the default sounding 
# latestRadarLoc = df.loc[4, 'Value'] # 'default' or 3-letter radar identifier to use for most recent radar image
# latestSoundLocs = df.loc[5, 'Value'].split(',') # 'default' or list of 4 3-letter identifiers for rawinsonde locations (latest launch time)
# fcstSoundLoc = df.loc[6, 'Value'] # 'default' 3-letter identifier for SBU-WRF forecast soundings

print('Generating the {} briefing using model data from the {}'.format(presentationType, str(modelList)[1:-1]))
print('   Using {} as the start time of the Day 0 period'.format(datetime.strftime(present_time, '%Y-%m-%d %H:%M')))
print('   Using {} as the initialization time for the model runs'.format(nearest6hr_shortStr))

##################################################
### POPULATE MODEL PRODUCT DIRECTORY
if presentationType=='morning':
    modelProducts = {'z500_vort':{'name':'500hPa Height, Vorticity', 'models':{'gfs', 'nam'}, 'scope':{'us'},
                                  'times':{0:[18,24,30], 1:[12,18,24,30], 2:[12,18,24,30]}},
                     'temp_adv_fgen_700':{'name':'700hPa TAdv, Fronto', 'models':{'gfs', 'nam'}, 'scope':{'us'},
                                          'times':{0:[18,24,30], 1:[12,18,24,30], 2:[12,18,24,30]}},
                     'ref_frzn':{'name':'dBZ, MSLP', 'models':{'gfs', 'nam'}, 'scope':{'us'},
                                 'times':{0:[18,24,30], 1:[12,18,24,30], 2:[12,18,24,30]}},
                     'refl_10cm':{'name':'dBZ, MSLP', 'models':{'wrfgfs', 'wrfnam'}, 'scope':{'eus'},
                                  'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
                                           1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '700_dBZfronto':{'name': '700hPa Fronto + dBZ', 'models':{'wrfgfs', 'wrfnam'},
                                      'scope':{'us', 'eus', 'ne'},
                                      'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
                                               1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     'xsect_MPV_dthetaEdz_fgen':{'name': '{} MPV, Fronto'.format(xSection),
                                                   'models':{'wrfgfs', 'wrfnam'}, 'scope':{'xsect'},
                                                   'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
                                                            1:[7,8,9,10,11,12,13,14,15,16,
                                                               17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
#                      'temps_700':{'name': '700hPa Temp, Height', 'models':{'wrfgfs', 'wrfnam'}, 'scope':{'eus'},
#                                   'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
#                                            1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
#                      'xsect_dBZfronto':{'name': '{} dBZ, Fronto'.format(xSection), 'models':{'wrfgfs', 'wrfnam'},
#                                         'scope':{'xsect'},
#                                         'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
#                                                  1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     'dbz1km': {'name': 'dBZ, MSLP', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '700hw': {'name': '700hPa Height, Vert. Vel.', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '500hv': {'name': '500hPa Height, Vort.', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '925hs': {'name': '925hPa Height, Wind', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     'z500_spag':{'name': '500hPa Height', 'models':{'gefs', 'naefs'}, 'scope':{'namer'},
                                  'times':{3:[12, 24], 4:[12, 24], 5:[12, 24]}},
                     'enslows':{'name': 'Pressure Centers', 'models':{'gefs', 'cmc'}, 'scope':{'namer'},
                                'times':{3:[12, 24], 4:[12, 24], 5:[12, 24]}},
                     'qpf_prob_025':{'name': 'Prob. 6-h QPF > 0.25', 'models':{'gefs'}, 'scope':{'namer'},
                                     'times':{3:[12, 24], 4:[12, 24], 5:[12, 24]}}
                    }
elif presentationType=='evening':
    modelProducts = {'z500_vort':{'name':'500hPa Height, Vorticity', 'models':{'gfs', 'nam'}, 'scope':{'us'},
                                  'times':{0:[24,30], 1:[12,18,24,30], 2:[12,18,24,30]}},
                     'temp_adv_fgen_700':{'name':'700hPa TAdv, Fronto', 'models':{'gfs', 'nam'}, 'scope':{'us'},
                                          'times':{0:[24,30], 1:[12,18,24,30], 2:[12,18,24,30]}},
                     'ref_frzn':{'name':'dBZ, MSLP', 'models':{'gfs', 'nam'}, 'scope':{'us'},
                                 'times':{0:[24,30], 1:[12,18,24,30], 2:[12,18,24,30]}},
                     'refl_10cm':{'name':'dBZ, MSLP', 'models':{'wrfgfs', 'wrfnam'}, 'scope':{'eus'},
                                  'times':{0:[24,25,26,27,28,29,30],
                                           1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '700_dBZfronto':{'name': '700hPa Fronto + dBZ', 'models':{'wrfgfs', 'wrfnam'},
                                      'scope':{'us', 'eus', 'ne'},
                                      'times':{0:[24,25,26,27,28,29,30],
                                               1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     'xsect_MPV_dthetaEdz_fgen':{'name': '{} MPV, Fronto'.format(xSection),
                                                   'models':{'wrfgfs', 'wrfnam'}, 'scope':{'xsect'},
                                                   'times':{0:[24,25,26,27,28,29,30],
                                                            1:[7,8,9,10,11,12,13,14,15,16,
                                                               17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
#                      'temps_700':{'name': '700hPa Temp, Height', 'models':{'wrfgfs', 'wrfnam'}, 'scope':{'eus'},
#                                   'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
#                                            1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
#                      'xsect_dBZfronto':{'name': '{} dBZ, Fronto'.format(xSection), 'models':{'wrfgfs', 'wrfnam'},
#                                         'scope':{'xsect'},
#                                         'times':{0:[18,19,20,21,22,23,24,25,26,27,28,29,30],
#                                                  1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     'dbz1km': {'name': 'dBZ, MSLP', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '700hw': {'name': '700hPa Height, Vert. Vel.', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '500hv': {'name': '500hPa Height, Vort.', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     '925hs': {'name': '925hPa Height, Wind', 'models':{'wrfuiuc'}, 'scope':{'mwus'},
                                'times':{0:[24,25,26,27,28,29,30],
                                         1:[7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,25,26,27,28,29,30]}},
                     'z500_spag':{'name': '500hPa Height', 'models':{'gefs', 'naefs'}, 'scope':{'namer'},
                                  'times':{3:[12, 24], 4:[12, 24], 5:[12, 24]}},
                     'enslows':{'name': 'Pressure Centers', 'models':{'gefs', 'cmc'}, 'scope':{'namer'},
                                'times':{3:[12, 24], 4:[12, 24], 5:[12, 24]}},
                     'qpf_prob_025':{'name': 'Prob. 6-h QPF > 0.25', 'models':{'gefs'}, 'scope':{'namer'},
                                     'times':{3:[12, 24], 4:[12, 24], 5:[12, 24]}}
                    }
# 	modelProducts = {'uv250':{'name':'250hPa_Wind_Height', 'models':{'gfs', 'nam'}, 'scope':{'conus'}, 'times':{0:[18,30], 1:[18, 30], 2:[18, 30], 3:[18, 42]}},\
# 		'z500_vort':{'name':'500hPa_Height_Vorticity', 'models':{'gfs', 'nam'}, 'scope':{'conus'}, 'times':{0:[18,30], 1:[18, 30], 2:[18, 30], 3:[18, 42]}},\
# 		'500hs':{'name':'500hPa_HeightSpaghetti', 'models':{'gfs'}, 'scope':{'namer'}, 'times':{3:[12, 36, 60, 84]}},\
# 		'midRH':{'name':'700-400hPa_RH_Wind_MSLP', 'models':{'gfs'}, 'scope':{'conus'}, 'times':{0:[18,30], 1:[18, 30], 2:[18, 30], 3:[18, 42]}},\
# 		'temp_adv_fgen_700':{'name':'700hPa_Temp_TAdv_Fronto_Wind', 'models':{'gfs', 'nam'}, 'scope':{'conus', 'eus'}, 'times':{0:[18,24,30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30], 3:[18, 42]}},\
# # 		'temp_adv_fgen_700':{'name':'700hPa_Temp_TAdv_Fronto_Wind', 'models':{'nam'}, 'scope':{'conus', 'eus'}, 'times':{0:[18,24,30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30]}},\
# 		'temp_adv_fgen_850':{'name':'850hPa_Temp_TAdv_Fronto_Wind', 'models':{'gfs', 'nam'}, 'scope':{'conus', 'eus'}, 'times':{0:[18,24,30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30], 3:[18, 42]}},\
# # 		'temp_adv_fgen_850':{'name':'850hPa_Temp_TAdv_Fronto_Wind', 'models':{'nam'}, 'scope':{'conus', 'eus'}, 'times':{0:[18,24,30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30]}},\
# 		'mslp_pcpn_frzn':{'name':'PrecipRate_MSLP_PrecipRate', 'models':{'gfs'}, 'scope':{'conus', 'eus'}, 'times':{0:[18, 24, 30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30], 3:[18, 42]}},\
# 		'ref_frzn':{'name':'dBZ_MSLP', 'models':{'nam', 'hrrr'}, 'scope':{'conus', 'eus'}, 'times':{0:[18, 24, 30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30]}},\
# 		'snbnd':{'name':'snowbandProb', 'models':{'hrrr'}, 'scope':{'ne'}, 'times':{0:[15, 18, 21, 24]}},\
# 		'ir':{'name':'SimulBrightTemp', 'models':{'hrrr'}, 'scope':{'eus'}, 'times':{0:[18, 24, 30]}},\
# 		'skewt':{'name':'ForecastSounding', 'models':{'gfs', 'nam'}, 'scope':{'sound'}, 'times':{0:[18, 24, 30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30]}},\
# 		'dBZtheta':{'name':'dBZ_Theta', 'models':{'gfs', 'nam'}, 'scope':{'xsect'}, 'times':{0:[18, 24, 30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30]}},\
# 		'mpvthes':{'name':'MPV_EquivSatTheta', 'models':{'gfs', 'nam'}, 'scope':{'xsect'}, 'times':{0:[18, 24, 30], 1:[12, 18, 24, 30], 2:[12, 18, 24, 30]}}
# 		} # JF: Add soundings
##################################################
### POPULATE REMOTE AND LOCAL PATHS FOR PLOTS
# Variable follows the pattern 'Product name': (remote_path [URL], local_path [filename only], variable string)
# First gather the observational plots

[wal_url, wal_datestr] = get_latest_radar('DOX') # Get timestamp info for the most recent radar image
[chs_url, chs_datestr] = get_latest_radar('CLX') # Get timestamp info for the most recent radar image
img_paths = {'z500_minus12_uair_us':
             ('/home/disk/funnel/impacts/archive/ops/upper_air/{}/ops.upper_air.{}.500mb.gif'.format(
                 past12hr_fullStr[0:8], past12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/upper_air/{}/ops.upper_air.{}.500mb.gif'.format(
                  past12hr_fullStr[0:8], past12hr_fullStr), '500-mb Heights',
              past12hr_shortStr, past12hr_locStr, past12hr_delta),
             'z500_current_uair_us':
             ('/home/disk/funnel/impacts/archive/ops/upper_air/{}/ops.upper_air.{}.500mb.gif'.format(
                 current12hr_fullStr[0:8], current12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/upper_air/{}/ops.upper_air.{}.500mb.gif'.format(
                  current12hr_fullStr[0:8], current12hr_fullStr), '500-mb Heights',
              current12hr_shortStr, current12hr_locStr, current12hr_delta),
             'anl_minus12_surf_atl': 
             ('/home/disk/funnel/impacts/archive/ops/sfc_anal/{}/ops.sfc_anal.{}.atlantic.gif'.format(
                 past12hr_fullStr[0:8], past12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/sfc_anal/{}/ops.sfc_anal.{}.atlantic.gif'.format(
                  past12hr_fullStr[0:8], past12hr_fullStr), 'Surface Analysis',
              past12hr_shortStr, past12hr_locStr, past12hr_delta),
             'anl_current_surf_atl':
             ('/home/disk/funnel/impacts/archive/ops/sfc_anal/{}/ops.sfc_anal.{}.atlantic.gif'.format(
                 current12hr_fullStr[0:8], current12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/sfc_anal/{}/ops.sfc_anal.{}.atlantic.gif'.format(
                  current12hr_fullStr[0:8], current12hr_fullStr), 'Surface Analysis',
              current12hr_shortStr, current12hr_locStr, current12hr_delta),
             'ir_minus12_conus':
             ('/home/disk/funnel/impacts/archive/ops/goes_east/{}/ops.goes_east.{}10.ir_4km.gif'.format(
                 past12hr_fullStr[0:8], past12hr_fullStr[0:10]),
              '/home/disk/funnel/impacts/archive/ops/goes_east/{}/ops.goes_east.{}10.ir_4km.gif'.format(
                 past12hr_fullStr[0:8], past12hr_fullStr[0:10]), 'Satellite IR',
              past12hr_shortStr, past12hr_locStr, past12hr_delta),
             'ir_current_conus':
             ('/home/disk/funnel/impacts/archive/ops/goes_east/{}/ops.goes_east.{}10.ir_4km.gif'.format(
                 nearest1hr_fullStr[0:8], nearest1hr_fullStr[0:10]),
              '/home/disk/funnel/impacts/archive/ops/goes_east/{}/ops.goes_east.{}10.ir_4km.gif'.format(
                 nearest1hr_fullStr[0:8], nearest1hr_fullStr[0:10]), 'Satellite IR',
              current12hr_shortStr, current12hr_locStr, current12hr_delta),
             'rad_minus12_conus': (['http://weather.rap.ucar.edu/radar/nids/images/N0R/KUSA/' + past12hr_fullStr[0:8] +
                                    '_' + past12hr_fullStr[8:] + '00BIG.png'][0], 'radar_past_conus.gif', 'CONUS Radar',
                                   past12hr_shortStr, past12hr_locStr, past12hr_delta),
             'rad_current_us': (['http://weather.rap.ucar.edu/radar/nids/images/N0R/KUSA/' + current_fullStr[0:8] +
                                 '_' + current_fullStr[8:10] + '0000BIG.png'][0], 'radar_conus.gif', 'CONUS Radar',
                                current_shortStr, current_locStr, current_delta),
             #'rad_minus12_conus':
             #('https://www2.mmm.ucar.edu/imagearchive1/RadarComposites/national/{}/national_{}.gif'.format(
                 #past12hr_fullStr[0:8], past12hr_fullStr), 'radar_past_conus.gif', 'CONUS Radar',
              #past12hr_shortStr, past12hr_locStr, past12hr_delta),
             #'rad_current_us': ('https://radar.weather.gov/ridge/Conus/RadarImg/latest_Small.gif', 'radar_conus.gif',
                                #'CONUS Radar', current_shortStr, current_locStr, current_delta),
             'rad_current_wal': ('{}{}_BREF1_color.png'.format(wal_url, wal_datestr), 'radar_wal.png', 'Dover Radar',
                                 current_shortStr, current_locStr, current_delta),
             'rad_current_chs': ('{}{}_BREF1_color.png'.format(chs_url, chs_datestr), 'radar_chs.png', 'Charleston Radar',
                                 current_shortStr, current_locStr, current_delta),
             'skewt_current_wal':
             ('/home/disk/funnel/impacts/archive/ops/skewt/{}/ops.skewt.{}.WAL.png'.format(
                 current12hr_fullStr[0:8], current12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/skewt/{}/ops.skewt.{}.WAL.png'.format(
                  current12hr_fullStr[0:8], current12hr_fullStr), 'Wallops Sounding',
              current12hr_shortStr, current12hr_locStr, current12hr_delta),
             'skewt_current_chs':
             ('/home/disk/funnel/impacts/archive/ops/skewt/{}/ops.skewt.{}.CHS.png'.format(
                 current12hr_fullStr[0:8], current12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/skewt/{}/ops.skewt.{}.CHS.png'.format(
                  current12hr_fullStr[0:8], current12hr_fullStr), 'Charleston Sounding',
              current12hr_shortStr, current12hr_locStr, current12hr_delta),
             'haz_current_ne': ('https://maps8.pivotalweather.com/maps/warnings/nwshaz.us_ne.png', 'hazards.png',
                                'Current Hazards', current_shortStr, current_locStr, current_delta),
             'snprob_day2_us':
             ('/home/disk/funnel/impacts/archive/ops/noaa/{}/ops.noaa.{}.day2_psnow_gt_04.gif'.format(
                 current12hr_fullStr[0:8], current12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/noaa/{}/ops.noaa.{}.day2_psnow_gt_04.gif'.format(
                  current12hr_fullStr[0:8], current12hr_fullStr), 'Snow Prob. > 4 in. (Day 2)',
              current12hr_shortStr, current12hr_locStr, current12hr_delta),
             'snprob_day3_us':
             ('/home/disk/funnel/impacts/archive/ops/noaa/{}/ops.noaa.{}.day3_psnow_gt_04.gif'.format(
                 current12hr_fullStr[0:8], current12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/noaa/{}/ops.noaa.{}.day3_psnow_gt_04.gif'.format(
                  current12hr_fullStr[0:8], current12hr_fullStr), 'Snow Prob. > 4 in. (Day 3)',
              current12hr_shortStr, current12hr_locStr, current12hr_delta),
             'lowtrks_us':
             ('/home/disk/funnel/impacts/archive/ops/noaa/{}/ops.noaa.{}.lowtrack.gif'.format(
                 nearest12hr_fullStr[0:8], nearest12hr_fullStr),
              '/home/disk/funnel/impacts/archive/ops/noaa/{}/ops.noaa.{}.lowtrack.gif'.format(
                  nearest12hr_fullStr[0:8], nearest12hr_fullStr), 'Low Track Forecast',
              nearest12hr_shortStr, nearest12hr_locStr, nearest12hr_delta)
            }

# img_paths = {'satIR_m1': ('https://cdn.star.nesdis.noaa.gov/GOES16/ABI/CONUS/14/' + dminus1_sat_datestr + '_GOES16-ABI-CONUS-14-2500x1500.jpg', 'DayM1_satIR.jpg', 'Satellite IR'), # GOES IR from previous day
# 	'satWV_m1': ('https://cdn.star.nesdis.noaa.gov/GOES16/ABI/CONUS/08/' + dminus1_sat_datestr + '_GOES16-ABI-CONUS-08-2500x1500.jpg', 'DayM1_satWV.jpg', 'Satellite WV'), # GOES WV from previous day
# 	'satIR_0': ('https://cdn.star.nesdis.noaa.gov/GOES16/ABI/CONUS/14/1250x750.jpg', 'Day0_satIR.jpg', 'Satellite IR'), # GOES IR closest to briefing time
# 	'satWV_0': ('https://cdn.star.nesdis.noaa.gov/GOES16/ABI/CONUS/08/1250x750.jpg', 'Day0_satWV.jpg', 'Satellite WV'), # GOES WV closest to briefing time
# 	'rad_nat': ('https://radar.weather.gov/ridge/Conus/RadarImg/latest.gif', 'radar_conus.gif', 'CONUS Radar'), # Latest radar mosaic of CONUS
# 	'rad_reg': ('https://radar.weather.gov/ridge/Conus/RadarImg/'+region_radar+'.gif', 'radar_'+region[2:]+'.gif', 'Regional Radar'), # Latest radar mosaic from the specified region
# 	'rad_local': (rad_url+rad_datestr+'_BREF1_color.png', 'radar_'+latestRadarLoc +'.png', latestRadarLoc+' Radar'), # Latest radar reflectivity from the selected site
# 	'ua': ('https://ocean.weather.gov/UA/East_coast.gif', 'surf_analysis.gif', 'Surface Analysis'), # Latest East Coast unified surface analysis
# 	'metar': ('http://weather.rap.ucar.edu/surface/' + metar_datestr + '_metars_' + region_metar + '.gif', 'metar.gif', 'Surface Observations'), # Recent surface METARs from the NE
# 	'skewt_1': ('http://impacts.atmos.washington.edu/archive/ops/skewt/'+sndg_datestr[:-4]+'/ops.skewt.'+sndg_datestr+'.'+latestSoundLocs[0]+'.png', 'skewt_'+latestSoundLocs[0]+'.png', latestSoundLocs[0]),
# 	'skewt_2': ('http://impacts.atmos.washington.edu/archive/ops/skewt/'+sndg_datestr[:-4]+'/ops.skewt.'+sndg_datestr+'.'+latestSoundLocs[1]+'.png', 'skewt_'+latestSoundLocs[1]+'.png', latestSoundLocs[1]),
# 	'skewt_3': ('http://impacts.atmos.washington.edu/archive/ops/skewt/'+sndg_datestr[:-4]+'/ops.skewt.'+sndg_datestr+'.'+latestSoundLocs[2]+'.png', 'skewt_'+latestSoundLocs[2]+'.png', latestSoundLocs[2]),
# 	'skewt_4': ('http://impacts.atmos.washington.edu/archive/ops/skewt/'+sndg_datestr[:-4]+'/ops.skewt.'+sndg_datestr+'.'+latestSoundLocs[3]+'.png', 'skewt_'+latestSoundLocs[3]+'.png', latestSoundLocs[3]),
# 	'haz': ('https://maps8.pivotalweather.com/maps/warnings/nwshaz.us_ne.png', 'hazards.png', 'Current Hazards'), # Current NE hazards map
# 	'snprob_0': ('https://origin.wpc.ncep.noaa.gov/pwpf_24hr/prb_24hsnow_ge04_' + snprob_datestr + 'f024.gif', 'Day0_snprob.gif', 'Prob. 24-h Snow > 4 in.'), # 24-hr forecast probability of snow greater than 4 inches - Day 0
# 	'snprob_1': ('https://origin.wpc.ncep.noaa.gov/pwpf_24hr/prb_24hsnow_ge04_' + snprob_datestr + 'f048.gif', 'Day1_snprob.gif', 'Prob. 24-h Snow > 4 in.'), # 24-hr forecast probability of snow greater than 4 inches - Day 1
# 	}

# Now gather the model plots
# Products to skip. Follows the tuple format: (model, product name)
skip = [('','')] 
print('\nInitialization complete.\nGathering URLs and assigning image paths for all products...')
for product in modelProducts.keys(): # Loop through available products
    name = modelProducts[product]['name']
    print('  Obtaining links for {}.'.format(name))
    
    for model in modelList: # Loop through desired models
        # Skip this model if it is not output for this product
        if (model, product) in skip:
            continue
        elif model not in modelProducts[product]['models']:
            continue

        for day in modelProducts[product]['times'].keys():
            # Check if current model has data out to the current day
            if model=='nam' and day>2:
                outsideRange = 1
            elif model=='hrrr' and day>0:
                outsideRange = 1
            else:
                outsideRange = 0

            if outsideRange==0: # Only gather URL and file paths if current model has data out to the current day
                i = 0
                for hour in modelProducts[product]['times'][day]: # Loop through the different forecast times
                    fcst_dt = present_time.replace(hour=0) + timedelta(days=day, hours=hour)
                    fcst_fullStr, fcst_shortStr, fcst_locStr, fcst_delta = dt2string(fcst_dt, present_time)
                    fhr_nearest6hr = int((fcst_dt - nearest6hr).total_seconds() / 3600)
                    fhr_nearest12hr = int((fcst_dt - nearest12hr).total_seconds() / 3600)
                    
#                     model_fullStr, model_shortStr, model_localStr, model_elapsed = dt2string(
#                         present_time.replace(hour=0) + timedelta(days=day, hours=hour), model_init_date)
#                     wrf_fullStr, wrf_shortStr, wrf_localStr, wrf_elapsed = dt2string(
#                         present_time.replace(hour=0) + timedelta(days=day, hours=hour), wrf_init_date)
#                     gefs_fullStr, gefs_shortStr, gefs_localStr, gefs_elapsed = dt2string(
#                         present_time.replace(hour=0) + timedelta(days=day, hours=hour), present_time)
#                     naefs_fullStr, naefs_shortStr, naefs_localStr, naefs_elapsed = dt2string(
#                         present_time.replace(hour=0) + timedelta(days=day, hours=hour), )
                    
#                     current_time = model_init_date.replace(hour=0) + timedelta(days=day, hours=hour)
#                     model_fcstTime = model_init_date.replace(hour=0) + timedelta(days=day, hours=hour)
#                     model_fcstHour = int((model_fcstTime - model_init_date).total_seconds() / 3600)
#                     wrf_fcstTime = wrf_init_date.replace(hour=0) + timedelta(days=day, hours=hour)
#                     wrf_fcstHour = int((wrf_fcstTime - wrf_init_date).total_seconds() / 3600)
#                     hrrr_fcstTime = hrrr_init_date.replace(hour=0) + timedelta(days=day, hours=hour) # for HRRR snowband product ONLY
#                     hrrr_fcstHour = int((hrrr_fcstTime - hrrr_init_date).total_seconds() / 3600)
#                     fcst_utcTime = datetime.strftime(present_time.replace(hour=0) +
#                                                      timedelta(days=day, hours=hour), '%HZ %m/%d')
#                     fcst_localTime = datetime.strftime(briefing_utcTime - timedelta(hours=5), '%-I%p %m/%d')
#                     fcst_elapsedTime = int((fcst_utcTime - present_time).total_seconds() / 3600)
                    
                    for scope in sorted(modelProducts[product]['scope']):
                        product_name = '{}_D{}H{}_{}_{}'.format(product, str(day), str(hour).zfill(2), model, scope)
                        if 'xsect' in scope: # Model cross section
                            remote_file = 'http://itpa.somas.stonybrook.edu/sbuwrf/{}/{}.{}_{}.d03.{}.gif'.format(
                                nearest12hr_fullStr[:10], model[3:].upper(), product, xSection_value, str(fhr_nearest12hr))
                            local_file = '{}_{}.gif'.format(product_name, xSection)
                            product_string = '{}-{} {}'.format(model[:3].upper(), model[3:].upper(),
                                                               modelProducts[product]['name'])
                        elif 'sound' in scope: # Forecast sounding
                            remote_file = 'http://itpa.somas.stonybrook.edu/sbuwrf/{}/{}.sound_{}.d03.{}.gif'.format(
                                nearest12hr_fullStr[:10], model[3:].upper(), fcstSoundLoc, str(fhr_nearest12hr))
                            local_file = '{}_{}.gif'.format(product_name, fcstSoundLoc)
                            product_string = '{}-{} {} Sounding'.format(model[:3].upper(), model[3:].upper(), fcstSoundLoc)
                        elif 'namer' in scope: # ensemble spaghetti plots
                            if model=='gefs':
                                if product=='z500_spag':
                                    remote_file = ['https://mag.ncep.noaa.gov/data/gefs-spag/' +
                                                   str(nearest12hr.hour).zfill(2) + '/gefs-spag_namer_' +
                                                   str(fhr_nearest12hr).zfill(3) + '_500_534_576_ht.gif'][0]
                                elif product=='qpf_prob_025':
                                    remote_file = ['https://mag.ncep.noaa.gov/data/gefs-mean-sprd/' +
                                                   str(nearest12hr.hour).zfill(2) + '/gefs-mean-sprd_namer_' +
                                                   str(fhr_nearest12hr).zfill(3) + '_prob_precip_0.25in.gif'][0]
                                elif product=='enslows':
                                    remote_file = ['https://tropicaltidbits.com/analysis/models/gfs-ens/' +
                                               nearest12hr_fullStr[:10] + '/gfs-ememb_lowlocs_us_' +
                                               str(int(fhr_nearest12hr/6)+1) + '.png'][0]
                            elif model=='naefs':
                                remote_file = ['https://weather.gc.ca/data/ensemble/images/' + nearest12hr_fullStr[:10] +
                                               '_054_E1_north@america_I_ENSEMBLE_spag@534_' +
                                               str(fhr_nearest12hr).zfill(3) + '.png'][0]
                            elif model=='cmc':
                                remote_file = ['https://tropicaltidbits.com/analysis/models/gem-ens/' +
                                               nearest12hr_fullStr[:10] + '/gem-ememb_lowlocs_us_' +
                                               str(int(fhr_nearest12hr/6)+1) + '.png'][0]
                            local_file = '{}.png'.format(product_name)
                            product_string = '{} {}'.format(model.upper(), modelProducts[product]['name'])
                        elif 'ne' in scope: # HRRR snowband product from the WPC
                            local_file = 'Day' + str(day) + '_' + product + str(chr(65 + i)) + '.gif'
                            remote_file = ['https://www.wpc.ncep.noaa.gov/snowbands/images/MTD_HRRRTLE_NE_snowobj_byhour_' +
                                           nearest12hr_fullStr[:10] + '_p1.00_t0.1_h' + str(fhr_nearest12hr).zfill(2) +
                                           '.png'][0]
                            product_string = 'HRRR Snowband Probability Product'
                        else: # Traditional plan view forecast graphics
                            if (model=='gfs') or (model=='nam'): # Tropical Tidbits graphis archived in the field catalog
                                
                                # === BLOCK IF USING IMAGES IN FIELD CATALOG === #
                                model_string = 'gfs_28km' if model=='gfs' else 'nam_12km'
                                remote_file = ['/home/disk/funnel/impacts/archive/model/' + model_string +
                                               '/' + nearest6hr_fullStr[:8] + '/model.' + model_string +
                                               '.' + nearest6hr_fullStr + '.' + str(fhr_nearest6hr).zfill(2) +
                                               '_' + product + '_' + scope + '.png'][0]
                                local_file = remote_file
                                # ========================= #
                                
                                '''
                                # === BLOCK IF USING IMAGES FROM TROPICAL TIDBITS === #
                                if model=='gfs':
                                    fhr_string = [str(int(fhr_nearest6hr/6)) if product=='ref_frzn' else
                                                  str(int(fhr_nearest6hr/6)+1)][0]
                                    remote_file = ['https://tropicaltidbits.com/analysis/models/gfs/' +
                                                   nearest6hr_fullStr[:10] + '/gfs_' +
                                                   product + '_us' + '_' + fhr_string + '.png'][0]
                                else:
                                    if fhr_nearest6hr<=36:
                                        fhr_string = [str(int(fhr_nearest6hr)) if product=='ref_frzn' else
                                                      str(int(fhr_nearest6hr)+1)][0]
                                    else:
                                        fhr_string = [str(int(36+(fhr_nearest6hr-36)/3)) if product=='ref_frzn' else
                                                      str(int(37+(fhr_nearest6hr-36)/3))][0]
                                    remote_file = ['https://tropicaltidbits.com/analysis/models/namconus/' +
                                                   nearest6hr_fullStr[:10] + '/namconus_' +
                                                   product + '_us' + '_' + fhr_string + '.png'][0]
                                local_file = '{}.png'.format(product_name)
                                # ========================= #
                                '''
                                product_string = '{} {}'.format(model.upper(), modelProducts[product]['name'])
                            elif (model=='wrfgfs') or (model=='wrfnam'):
                                remote_file = 'http://itpa.somas.stonybrook.edu/sbuwrf/{}/{}.{}.d02.{}.gif'.format(
                                    nearest12hr_fullStr[:10], model[3:].upper(), product, str(fhr_nearest12hr))
                                local_file = '{}.gif'.format(product_name)
                                product_string = '{} {}'.format(model.upper(), modelProducts[product]['name'])
                        img_paths[product_name] = (remote_file, local_file, product_string,
                                                   fcst_shortStr, fcst_locStr, fcst_delta)
                        print(remote_file)
# 						else: # Plan view spanning CONUS or the specified region
# 							local_file = 'Day' + str(day) + '_' + product + str(chr(65 + i)) + '_' + model + '_' + scope + '.png'
# 							product_string = model.upper() + ' ' + modelProducts[product]['name']
							
# 							if scope=='conus':
# 								scope = 'us'
							
# 							if model=='nam': # Slightly different URL scheme for NAM
# 								if model_fcstHour<=36:
# 									tstep = model_fcstHour + 1
# 								else:
# 									tstep = int(38+(model_fcstHour-39)/3)
									
# 								if product=='ref_frzn':
# 									tstep = tstep - 1
# 								remote_file = 'https://www.tropicaltidbits.com/analysis/models/namconus/' + model_init_string + '/namconus_' + product + '_' + scope + '_' + str(tstep) + '.png'
# 							elif model=='gfs':
# 								tstep = int(model_fcstHour / 6 + 1)
# 								if product=='mslp_pcpn_frzn':
# 									tstep = tstep - 1
# 								remote_file = 'https://www.tropicaltidbits.com/analysis/models/' + model + '/' + model_init_string + '/' + model + '_' + product + '_' + scope + '_' + str(tstep) + '.png'
# 							else: # HRRR
# 								tstep = hrrr_fcstHour + 1
# 								if product=='ref_frzn':
# 									tstep = tstep - 1
# 								remote_file = 'https://www.tropicaltidbits.com/analysis/models/' + model + '/' + hrrr_init_string + '/' + model + '_' + product + '_' + scope + '_' + str(tstep) + '.png'
# 						img_paths[product_name] = (remote_file, local_file, product_string)
##################################################
### BUILD POWERPOINT FILE
# Different default possibilities for the slide layout
layout = {'Title Slide' : 0, 'Bullet Slide' : 1, 'Segue' : 2, 'Side By Side' : 3,
          'Title Alone' : 5, 'Blank Slide' : 6, 'Picture with Caption' : 8
         }
         
def build_presentation(nearest6hr, present_time):
    '''
    This function builds the forecasting briefling PowerPoint template.
    Inputs:
        model_init_date: the initialization time for the model graphics
        present_time: the time the briefing is given
    '''
    print('\nBuilding the PowerPoint presentation...\n  Images will be saved to {}.'.format(downloadPath))
    
    # Switch to the presentation directory
    basedir = os.getcwd()
    os.chdir(downloadPath)

    # Make a new image subdirectory for this presentation
    if os.path.exists('./{}'.format(datetime.strftime(present_time, '%Y%m%d%H'))) is False:
        os.mkdir('./{}'.format(datetime.strftime(present_time, '%Y%m%d%H')))
    os.chdir('./{}'.format(datetime.strftime(present_time, '%Y%m%d%H')))
    
    if briefingUpdate=='True':
        if os.path.exists('./updated') is False:
            os.mkdir('./updated')
        os.chdir('./updated')
    
    # Make the presentation object
    prs = Presentation()

    # Make the title slide and choose the layout
    title_slide_layout = prs.slide_layouts[layout['Title Slide']]
    title_slide = prs.slides.add_slide(title_slide_layout) # Add it to the presentation
    if presentationType == 'morning':
        title_slide.shapes.title.text = 'Morning Weather Briefing'
    elif presentationType == 'evening':
        title_slide.shapes.title.text = 'Evening Weather Briefing'

    # Subtitle is a "placeholder" object
    if presentationType == 'morning':
        title_slide.placeholders[1].text = "{} 1400 UTC\nForecaster Name".format(datetime.strftime(present_time, '%d %b %Y'))
    elif presentationType == 'evening':
        title_slide.placeholders[1].text = "{} 2300 UTC\nForecaster Name".format(datetime.strftime(present_time, '%d %b %Y'))
        
    ### START OF THE BRIEFING POWERPOINT TEMPLATE
    prs = full_summary(prs, 'Briefing Overview', -1)
    if show_weather=='True':
        # Day -1 slides
        print('\n  Making Day -1 slides')
        prs = bumper_slide(prs, 'Past {} Hours'.format(str(-1*past12hr_delta)), -1, past12hr, present_time)
        prs = full_summary(prs, 'Summary of Past {} Hours'.format(str(-1*past12hr_delta)), -1)
        prs = four_panel_image(prs, ['rad_minus12_conus', 'ir_minus12_conus',
                                     'z500_minus12_uair_us', 'anl_minus12_surf_atl'], -1)

        # Current weather slides
        print('\n  Making Current Weather slides')
        prs = four_panel_image(prs, ['rad_current_us', 'ir_current_conus',
                                     'z500_current_uair_us', 'anl_current_surf_atl'], 0)
        prs = full_slide_image(prs, 'haz_current_ne', 0, datetime.utcnow(), width=8, height=6.9)
        prs = two_panel_image(prs, ['rad_current_wal', 'skewt_current_wal'], 0, [datetime.utcnow(), datetime.utcnow()],
                              title='Wallops Conditions & TAF')
        prs = two_panel_image(prs, ['rad_current_chs', 'skewt_current_chs'], 0, [datetime.utcnow(), datetime.utcnow()],
                              title='Hunter Conditions & TAF')

    if show_shortTerm=='True':
        prs = bumper_slide(prs, 'Synoptic Forecast', [0, 1, 2], datetime(utcnow.year, utcnow.month, utcnow.day, 14, 0),
                           datetime(utcnow.year, utcnow.month, utcnow.day+3, 6, 0))
        prs = timeline_slide(prs, [0, 1, 2], presentationType, products=['snprob_day2_us', 'lowtrks_us'])

        # Day 0 slides
        print('\n Making Day 0 slides')
        day = 0
        if presentationType=='morning':
            hourList = [18, 24, 30]
        elif presentationType=='evening':
            hourList = [24, 30]
            
        for hour in hourList:
            product1 = 'z500_vort_D0H{}_gfs_us'.format(str(hour).zfill(2))
            product2 = 'temp_adv_fgen_700_D0H{}_gfs_us'.format(str(hour).zfill(2))
            product3 = 'ref_frzn_D0H{}_gfs_us'.format(str(hour).zfill(2))
            product4 = 'z500_vort_D0H{}_nam_us'.format(str(hour).zfill(2))
            product5 = 'temp_adv_fgen_700_D0H{}_nam_us'.format(str(hour).zfill(2))
            product6 = 'ref_frzn_D0H{}_nam_us'.format(str(hour).zfill(2))
            prs = six_panel_image(prs, [product1, product2, product3, product4, product5, product6], 0)

        # Day 1 and 2 slides
        for day in [1, 2]:
            print('\n Making Day {} slides'.format(str(day)))
            for hour in [12, 18, 24, 30]:
                product1 = 'z500_vort_D{}H{}_gfs_us'.format(str(day), str(hour).zfill(2))
                product2 = 'temp_adv_fgen_700_D{}H{}_gfs_us'.format(str(day), str(hour).zfill(2))
                product3 = 'ref_frzn_D{}H{}_gfs_us'.format(str(day), str(hour).zfill(2))
                product4 = 'z500_vort_D{}H{}_nam_us'.format(str(day), str(hour).zfill(2))
                product5 = 'temp_adv_fgen_700_D{}H{}_nam_us'.format(str(day), str(hour).zfill(2))
                product6 = 'ref_frzn_D{}H{}_nam_us'.format(str(day), str(hour).zfill(2))
                prs = six_panel_image(prs, [product1, product2, product3, product4, product5, product6], day)

        # Day 0-2 Summary
        prs = timeline_slide(prs, [0, 1, 2], presentationType, products=['snprob_day2_us', 'lowtrks_us'])
        if region=='usne':
            get_gpm_overpasses(nearest1hr, region='east')
        else:
            get_gpm_overpasses(nearest1hr, region='west')
        if 'overpass1' in img_paths.keys(): # at least one GPM overpass was found...make a slide
            prs = six_panel_image(prs, ['overpass1', 'overpass2', 'overpass3', 'overpass4', 'overpass5', 'overpass6'],
                                  [0, 1, 2], plotTimebar=False)
        prs = objectives_slide(prs, 'Day 0-2 Summary', [0, 1, 2])
        
    if show_detailed=='True':
        # Detailed forecast slides
        prs = bumper_slide(prs, 'Detailed Forecast', [0, 1], datetime(utcnow.year, utcnow.month, utcnow.day, 14, 0),
                           datetime(utcnow.year, utcnow.month, utcnow.day+2, 6, 0))
        day = 0
        if presentationType=='morning':
            hourList = range(18, 31)
        elif presentationType=='evening':
            hourList = range(24, 31)

        for hour in hourList:
            if (region=='usne') and ('wrfgfs' in modelList) and ('wrfnam' in modelList):
                product1 = 'refl_10cm_D0H{}_wrfgfs_eus'.format(str(hour).zfill(2))
                product2 = '700_dBZfronto_D0H{}_wrfgfs_eus'.format(str(hour).zfill(2))
                product3 = 'xsect_MPV_dthetaEdz_fgen_D0H{}_wrfgfs_xsect'.format(str(hour).zfill(2))
                product4 = 'refl_10cm_D0H{}_wrfnam_eus'.format(str(hour).zfill(2))
                product5 = '700_dBZfronto_D0H{}_wrfnam_eus'.format(str(hour).zfill(2))
                product6 = 'xsect_MPV_dthetaEdz_fgen_D0H{}_wrfnam_xsect'.format(str(hour).zfill(2))
                prs = six_panel_image(prs, [product1, product2, product3, product4, product5, product6], 0)
            elif (region=='usmw') and ('wrfuiuc' in modelList):
                product1 = 'dbz1km_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                product2 = '500hv_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                product3 = '700hw_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                product4 = '925hs_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                prs = four_panel_image(prs, [product1, product2, product3, product4], 0)

        day = 1
        for hour in range(7, 31):
            if (region=='usne') and ('wrfgfs' in modelList) and ('wrfnam' in modelList):
                product1 = 'refl_10cm_D1H{}_wrfgfs_eus'.format(str(hour).zfill(2))
                product2 = '700_dBZfronto_D1H{}_wrfgfs_eus'.format(str(hour).zfill(2))
                product3 = 'xsect_MPV_dthetaEdz_fgen_D1H{}_wrfgfs_xsect'.format(str(hour).zfill(2))
                product4 = 'refl_10cm_D1H{}_wrfnam_eus'.format(str(hour).zfill(2))
                product5 = '700_dBZfronto_D1H{}_wrfnam_eus'.format(str(hour).zfill(2))
                product6 = 'xsect_MPV_dthetaEdz_fgen_D1H{}_wrfnam_xsect'.format(str(hour).zfill(2))
                prs = six_panel_image(prs, [product1, product2, product3, product4, product5, product6], 1)
            elif (region=='usmw') and ('wrfuiuc' in modelList):
                product1 = 'dbz1km_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                product2 = '500hv_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                product3 = '700hw_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                product4 = '925hs_D0H{}_wrfuiuc_mwus'.format(str(hour).zfill(2))
                prs = four_panel_image(prs, [product1, product2, product3, product4], 1)

        prs = objectives_slide(prs, 'Detailed Forecast Summary', [0, 1])

    if show_longTerm=='True':
        prs = bumper_slide(prs, 'Long-term Forecast', 3, datetime(utcnow.year, utcnow.month, utcnow.day+3, 6, 0),
                           datetime(utcnow.year, utcnow.month, utcnow.day+6, 6, 0))
        prs = timeline_slide(prs, 3, presentationType, products=None)
        #prs = two_panel_image(prs, ['snprob_day3_us', 'lowtrks_us'], 3, [datetime.utcnow(), datetime.utcnow()],
                              #title='Day 3 Overview')

        # Day 3+ slides
        print('\n  Making Day 3+ slides')
        for day in [3, 4, 5]:
            for hour in [12, 24]:
                product1 = 'z500_spag_D{}H{}_gefs_namer'.format(str(day), str(hour).zfill(2))
                product2 = 'enslows_D{}H{}_gefs_namer'.format(str(day), str(hour).zfill(2))
                #product2 = 'z500_spag_D{}H{}_naefs_namer'.format(str(day), str(hour).zfill(2))
                product3 = 'qpf_prob_025_D{}H{}_gefs_namer'.format(str(day), str(hour).zfill(2))
                product4 = 'enslows_D{}H{}_cmc_namer'.format(str(day), str(hour).zfill(2))
                prs = four_panel_image(prs, [product1, product2, product3, product4], 3)

        prs = objectives_slide(prs, 'Snow Plumes', 3)

        # Day 3+ Summary
        prs = objectives_slide(prs, 'Day 3+ Summary', day)

    # Save the presentation
    if briefingUpdate=='True':
        prs.save('{}/forecast.wxbriefing.{}.{}_new.pptx'.format(
            presentationPath, datetime.strftime(present_time, '%Y%m%d'), briefingString))
    else:
        prs.save('{}/forecast.wxbriefing.{}.{}.pptx'.format(
            presentationPath, datetime.strftime(present_time, '%Y%m%d'), briefingString))
    
if __name__ == '__main__':
    build_presentation(nearest6hr, present_time)
